######-----Import Dash-----#####
import dash
from dash import dcc
from dash import html, callback
import dash_bootstrap_components as dbc
from dash.dependencies import Input, Output
from dash.exceptions import PreventUpdate

from pptx import Presentation
from pptx.util import Cm, Pt

import os

import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.io as pio
from plotly.subplots import make_subplots

from datetime import date
from dateutil.relativedelta import relativedelta

dash.register_page(__name__, name='Download')

##-----page

## -----LAYOUT-----
layout = html.Div([
                html.Br(),
                
                html.Div([
                    html.Button('Download Slides', 
                                      id='download_button', 
                                      n_clicks=0, 
                                      style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana','display': 'inline-block'}), 
                    dcc.Download(id='download'),
                    ]),

                html.Br(),
                html.P(id='download_date', style={'fontSize': 15, 'color':'#2a3f5f','font-family':'Verdana'}),

    html.Footer('ABL',
            style={'textAlign': 'center', 
                   'fontSize': 20, 
                   'background-color':'#2a3f5f',
                   'color':'white',
                   'position': 'fixed',
                    'bottom': '0',
                    'width': '100%'})

    ], style={
        'paddingLeft':'10px',
        'paddingRight':'10px',
    })

#### Callback Auto Update Chart & Data

@callback(
    [Output('download', 'data'),
     Output('download_date', 'children')],
    [Input('download_button', 'n_clicks'),
     Input('store', 'data')]
)
def update_charts(n_clicks, data):
    if n_clicks == 0:
        raise PreventUpdate
    ######################
    # Data
    ######################
    dat_borneo = pd.DataFrame(data['Bulk Borneo'])
    dat_celebes = pd.DataFrame(data['Bulk Celebes'])
    dat_sumatra = pd.DataFrame(data['Bulk Sumatra'])
    dat_java = pd.DataFrame(data['Bulk Java'])
    dat_dewata = pd.DataFrame(data['Bulk Dewata'])
    dat_karimun = pd.DataFrame(data['Bulk Karimun'])
    dat_of1 = pd.DataFrame(data['Ocean Flow 1'])
    dat_natuna = pd.DataFrame(data['Bulk Natuna'])
    dat_sumba = pd.DataFrame(data['Bulk Sumba'])
    dat_derawan = pd.DataFrame(data['Bulk Derawan'])
    dat_greencalypso = pd.DataFrame(data['Green Calypso'])
    dat_putrialysha = pd.DataFrame(data['Putri Alysha'])
    dat_bunaken = pd.DataFrame(data['Bulk Bunaken'])
    dat_svii = pd.DataFrame(data['STRAITS VENTURE II'])

    # Chart
    volume_sumatra = plot_volume(dat_sumatra, "Volume Bulk Sumatra")
    volume_dewata = plot_volume(dat_dewata, "Volume Bulk Dewata")
    volume_karimun = plot_volume(dat_karimun, "Volume Bulk Karimun")
    volume_derawan = plot_volume(dat_derawan, "Volume Bulk Derawan")
    volume_of1 = plot_volume(dat_of1, "Volume Ocean Flow 1")
    volume_sumba = plot_volume(dat_sumba, "Volume Bulk Sumba")
    volume_java = plot_volume(dat_java, "Volume Bulk Java")
    volume_natuna = plot_volume(dat_natuna, "Volume Bulk Natuna")
    volume_celebes = plot_volume(dat_celebes, "Volume Bulk Celebes")
    volume_borneo = plot_volume(dat_borneo, "Volume Bulk Borneo")
    volume_greencalypso = plot_volume(dat_greencalypso, "Volume Green Calypso")
    volume_putrialysha = plot_volume(dat_putrialysha, "Volume Putri Alysha")
    volume_bunaken = plot_volume(dat_bunaken, "Volume Bulk Bunaken")
    volume_svii = plot_volume(dat_svii, "Volume STRAITS VENTURE II")

    nlr_sumatra = plot_nlr(dat_sumatra, "NLR Bulk Sumatra")
    nlr_dewata = plot_nlr(dat_dewata, "NLR Bulk Dewata")
    nlr_karimun = plot_nlr(dat_karimun, "NLR Bulk Karimun")
    nlr_derawan = plot_nlr(dat_derawan, "NLR Bulk Derawan")
    nlr_of1 = plot_nlr(dat_of1, "NLR Ocean Flow 1")
    nlr_sumba = plot_nlr(dat_sumba, "NLR Bulk Sumba")
    nlr_java = plot_nlr(dat_java, "NLR Bulk Java")
    nlr_natuna = plot_nlr(dat_natuna, "NLR Bulk Natuna")
    nlr_celebes = plot_nlr(dat_celebes, "NLR Bulk Celebes")
    nlr_borneo = plot_nlr(dat_borneo, "NLR Bulk Borneo")
    nlr_greencalypso = plot_nlr(dat_greencalypso, "NLR Green Calypso")
    nlr_putrialysha = plot_nlr(dat_putrialysha, "NLR Putri Alysha")
    nlr_bunaken = plot_nlr(dat_bunaken, "NLR Bulk Bunaken")
    nlr_svii = plot_nlr(dat_svii, "NLR STRAITS VENTURE II")

    glr_sumatra = plot_glr(dat_sumatra, "GLR Bulk Sumatra")
    glr_dewata = plot_glr(dat_dewata, "GLR Bulk Dewata")
    glr_karimun = plot_glr(dat_karimun, "GLR Bulk Karimun")
    glr_derawan = plot_glr(dat_derawan, "GLR Bulk Derawan")
    glr_of1 = plot_glr(dat_of1, "GLR Ocean Flow 1")
    glr_sumba = plot_glr(dat_sumba, "GLR Bulk Sumba")
    glr_java = plot_glr(dat_java, "GLR Bulk Java")
    glr_natuna = plot_glr(dat_natuna, "GLR Bulk Natuna")
    glr_celebes = plot_glr(dat_celebes, "GLR Bulk Celebes")
    glr_borneo = plot_glr(dat_borneo, "GLR Bulk Borneo")
    glr_greencalypso = plot_glr(dat_greencalypso, "GLR Green Calypso")
    glr_putrialysha = plot_glr(dat_putrialysha, "GLR Putri Alysha")
    glr_bunaken = plot_glr(dat_bunaken, "GLR Bulk Bunaken")
    glr_svii = plot_glr(dat_svii, "GLR STRAITS VENTURE II")

    fr_sumatra = plot_fr(dat_sumatra, "Fuel Ratio Bulk Sumatra", 0.24)
    fr_dewata = plot_fr(dat_dewata, "Fuel Ratio Bulk Dewata", 0.29)
    fr_karimun = plot_fr(dat_karimun, "Fuel Ratio Bulk Karimun", 0.25)
    fr_derawan = plot_fr(dat_derawan, "Fuel Ratio Bulk Derawan", 0.33)
    fr_of1 = plot_fr(dat_of1, "Fuel Ratio Ocean Flow 1", 0.32)
    fr_sumba = plot_fr(dat_sumba, "Fuel Ratio Bulk Sumba", 0)
    fr_java = plot_fr(dat_java, "Fuel Ratio Bulk Java", 0.27)
    fr_natuna = plot_fr(dat_natuna, "Fuel Ratio Bulk Natuna", 0.40)
    fr_celebes = plot_fr(dat_celebes, "Fuel Ratio Bulk Celebes", 0.37)
    fr_borneo = plot_fr(dat_borneo, "Fuel Ratio Bulk Borneo", 0.24)
    fr_greencalypso = plot_fr(dat_greencalypso, "Fuel Ratio Green Calypso", 0)
    fr_putrialysha = plot_fr(dat_putrialysha, "Fuel Ratio Putri Alysha", 0)
    fr_bunaken = plot_fr(dat_bunaken, "Fuel Ratio Bulk Bunaken", 0)
    fr_svii = plot_fr(dat_svii, "Fuel Ratio STRAITS VENTURE II", 0)

    nlr_type_sumatra = plot_nlr_type2(dat_sumatra, "NLR Bulk Sumatra", 52000, 41673)
    nlr_type_dewata = plot_nlr_type2(dat_dewata, "NLR Bulk Dewata", 36815, 44482)
    nlr_type_karimun = plot_nlr_type2(dat_karimun, "NLR Bulk Karimun", 46000, 26878)
    nlr_type_derawan = plot_nlr_type(dat_derawan, "NLR Bulk Derawan", 46000)
    nlr_type_of1 = plot_nlr_type(dat_of1, "NLR Ocean Flow 1", 36000)
    nlr_type_sumba = plot_nlr_type(dat_sumba, "NLR Bulk Sumba", 0)
    nlr_type_java = plot_nlr_type(dat_java, "NLR Bulk Java", 46000)
    nlr_type_natuna = plot_nlr_type(dat_natuna, "NLR Bulk Natuna", 18400)
    nlr_type_celebes = plot_nlr_type2(dat_celebes, "NLR Bulk Celebes", 46000, 39783)
    nlr_type_borneo = plot_nlr_type(dat_borneo, "NLR Bulk Borneo", 25000)
    nlr_type_greencalypso = plot_nlr_type(dat_greencalypso, "NLR Green Calypso", 0)
    nlr_type_putrialysha = plot_nlr_type(dat_putrialysha, "NLR Putri Alysha", 0)
    nlr_type_bunaken = plot_nlr_type(dat_bunaken, "NLR Bulk Bunaken", 0)
    nlr_type_svii = plot_nlr_type(dat_svii, "NLR STRAITS VENTURE II", 0)

    glr_type_sumatra = plot_glr_type(dat_sumatra, "GLR Bulk Sumatra")
    glr_type_dewata = plot_glr_type(dat_dewata, "GLR Bulk Dewata")
    glr_type_karimun = plot_glr_type(dat_karimun, "GLR Bulk Karimun")
    glr_type_derawan = plot_glr_type(dat_derawan, "GLR Bulk Derawan")
    glr_type_of1 = plot_glr_type(dat_of1, "GLR Ocean Flow 1")
    glr_type_sumba = plot_glr_type(dat_sumba, "GLR Bulk Sumba")
    glr_type_java = plot_glr_type(dat_java, "GLR Bulk Java")
    glr_type_natuna = plot_glr_type(dat_natuna, "GLR Bulk Natuna")
    glr_type_celebes = plot_glr_type(dat_celebes, "GLR Bulk Celebes")
    glr_type_borneo = plot_glr_type(dat_borneo, "GLR Bulk Borneo")
    glr_type_greencalypso = plot_glr_type(dat_greencalypso, "GLR Green Calypso")
    glr_type_putrialysha = plot_glr_type(dat_putrialysha, "GLR Putri Alysha")
    glr_type_bunaken = plot_glr_type(dat_bunaken, "GLR Bulk Bunaken")
    glr_type_svii = plot_glr_type(dat_svii, "GLR STRAITS VENTURE II")

    # Save img
    pio.write_image(volume_sumatra, 'img/Volume Bulk Sumatra.png')
    pio.write_image(volume_dewata, 'img/Volume Bulk Dewata.png')
    pio.write_image(volume_karimun, 'img/Volume Bulk Karimun.png')
    pio.write_image(volume_derawan, 'img/Volume Bulk Derawan.png')
    pio.write_image(volume_of1, 'img/Volume Ocean Flow 1.png')
    pio.write_image(volume_sumba, 'img/Volume Bulk Sumba.png')
    pio.write_image(volume_java, 'img/Volume Bulk Java.png')
    pio.write_image(volume_natuna, 'img/Volume Bulk Natuna.png')
    pio.write_image(volume_celebes, 'img/Volume Bulk Celebes.png')
    pio.write_image(volume_borneo, 'img/Volume Bulk Borneo.png')
    pio.write_image(volume_greencalypso, 'img/Volume Green Calypso.png')
    pio.write_image(volume_putrialysha, 'img/Volume Putri Alysha.png')
    pio.write_image(volume_bunaken, 'img/Volume Bulk Bunaken.png')
    pio.write_image(volume_svii, 'img/Volume STRAITS VENTURE II.png')

    pio.write_image(nlr_sumatra, 'img/NLR Bulk Sumatra.png')
    pio.write_image(nlr_dewata, 'img/NLR Bulk Dewata.png')
    pio.write_image(nlr_karimun, 'img/NLR Bulk Karimun.png')
    pio.write_image(nlr_derawan, 'img/NLR Bulk Derawan.png')
    pio.write_image(nlr_of1, 'img/NLR Ocean Flow 1.png')
    pio.write_image(nlr_sumba, 'img/NLR Bulk Sumba.png')
    pio.write_image(nlr_java, 'img/NLR Bulk Java.png')
    pio.write_image(nlr_natuna, 'img/NLR Bulk Natuna.png')
    pio.write_image(nlr_celebes, 'img/NLR Bulk Celebes.png')
    pio.write_image(nlr_borneo, 'img/NLR Bulk Borneo.png')
    pio.write_image(nlr_greencalypso, 'img/NLR Green Calypso.png')
    pio.write_image(nlr_putrialysha, 'img/NLR Putri Alysha.png')
    pio.write_image(nlr_bunaken, 'img/NLR Bulk Bunaken.png')
    pio.write_image(nlr_svii, 'img/NLR STRAITS VENTURE II.png')

    pio.write_image(glr_sumatra, 'img/GLR Bulk Sumatra.png')
    pio.write_image(glr_dewata, 'img/GLR Bulk Dewata.png')
    pio.write_image(glr_karimun, 'img/GLR Bulk Karimun.png')
    pio.write_image(glr_derawan, 'img/GLR Bulk Derawan.png')
    pio.write_image(glr_of1, 'img/GLR Ocean Flow 1.png')
    pio.write_image(glr_sumba, 'img/GLR Bulk Sumba.png')
    pio.write_image(glr_java, 'img/GLR Bulk Java.png')
    pio.write_image(glr_natuna, 'img/GLR Bulk Natuna.png')
    pio.write_image(glr_celebes, 'img/GLR Bulk Celebes.png')
    pio.write_image(glr_borneo, 'img/GLR Bulk Borneo.png')
    pio.write_image(glr_greencalypso, 'img/GLR Green Calypso.png')
    pio.write_image(glr_putrialysha, 'img/GLR Putri Alysha.png')
    pio.write_image(glr_bunaken, 'img/GLR Bulk Bunaken.png')
    pio.write_image(glr_svii, 'img/GLR STRAITS VENTURE II.png')

    pio.write_image(fr_sumatra, 'img/Fuel Ratio Bulk Sumatra.png')
    pio.write_image(fr_dewata, 'img/Fuel Ratio Bulk Dewata.png')
    pio.write_image(fr_karimun, 'img/Fuel Ratio Bulk Karimun.png')
    pio.write_image(fr_derawan, 'img/Fuel Ratio Bulk Derawan.png')
    pio.write_image(fr_of1, 'img/Fuel Ratio Ocean Flow 1.png')
    pio.write_image(fr_sumba, 'img/Fuel Ratio Bulk Sumba.png')
    pio.write_image(fr_java, 'img/Fuel Ratio Bulk Java.png')
    pio.write_image(fr_natuna, 'img/Fuel Ratio Bulk Natuna.png')
    pio.write_image(fr_celebes, 'img/Fuel Ratio Bulk Celebes.png')
    pio.write_image(fr_borneo, 'img/Fuel Ratio Bulk Borneo.png')
    pio.write_image(fr_greencalypso, 'img/Fuel Ratio Green Calypso.png')
    pio.write_image(fr_putrialysha, 'img/Fuel Ratio Putri Alysha.png')
    pio.write_image(fr_bunaken, 'img/Fuel Ratio Bulk Bunaken.png')
    pio.write_image(fr_svii, 'img/Fuel Ratio STRAITS VENTURE II.png')

    pio.write_image(nlr_type_sumatra, 'img/NLR_type Bulk Sumatra.png')
    pio.write_image(nlr_type_dewata, 'img/NLR_type Bulk Dewata.png')
    pio.write_image(nlr_type_karimun, 'img/NLR_type Bulk Karimun.png')
    pio.write_image(nlr_type_derawan, 'img/NLR_type Bulk Derawan.png')
    pio.write_image(nlr_type_of1, 'img/NLR_type Ocean Flow 1.png')
    pio.write_image(nlr_type_sumba, 'img/NLR_type Bulk Sumba.png')
    pio.write_image(nlr_type_java, 'img/NLR_type Bulk Java.png')
    pio.write_image(nlr_type_natuna, 'img/NLR_type Bulk Natuna.png')
    pio.write_image(nlr_type_celebes, 'img/NLR_type Bulk Celebes.png')
    pio.write_image(nlr_type_borneo, 'img/NLR_type Bulk Borneo.png')
    pio.write_image(nlr_type_greencalypso, 'img/NLR_type Green Calypso.png')
    pio.write_image(nlr_type_putrialysha, 'img/NLR_type Putri Alysha.png')
    pio.write_image(nlr_type_bunaken, 'img/NLR_type Bulk Bunaken.png')
    pio.write_image(nlr_type_svii, 'img/NLR_type STRAITS VENTURE II.png')

    pio.write_image(glr_type_sumatra, 'img/GLR_type Bulk Sumatra.png')
    pio.write_image(glr_type_dewata, 'img/GLR_type Bulk Dewata.png')
    pio.write_image(glr_type_karimun, 'img/GLR_type Bulk Karimun.png')
    pio.write_image(glr_type_derawan, 'img/GLR_type Bulk Derawan.png')
    pio.write_image(glr_type_of1, 'img/GLR_type Ocean Flow 1.png')
    pio.write_image(glr_type_sumba, 'img/GLR_type Bulk Sumba.png')
    pio.write_image(glr_type_java, 'img/GLR_type Bulk Java.png')
    pio.write_image(glr_type_natuna, 'img/GLR_type Bulk Natuna.png')
    pio.write_image(glr_type_celebes, 'img/GLR_type Bulk Celebes.png')
    pio.write_image(glr_type_borneo, 'img/GLR_type Bulk Borneo.png') 
    pio.write_image(glr_type_greencalypso, 'img/GLR_type Green Calypso.png')   
    pio.write_image(glr_type_putrialysha, 'img/GLR_type Putri Alysha.png') 
    pio.write_image(glr_type_bunaken, 'img/GLR_type Bulk Bunaken.png')
    pio.write_image(glr_type_svii, 'img/GLR_type STRAITS VENTURE II.png') 

    path = os.getcwd()
    download_date = html.Strong('Downloaded slides as per ' + str(dat_sumatra['Month'].iloc[-2]))
    ###################### PPTX ######################
    def to_pptx(bytes_io):
        # Create presentation
        pptx = path + '//' + 'slide_master.pptx'
        prs = Presentation(pptx)

        # define slidelayouts 
        slide = prs.slides.add_slide(prs.slide_layouts[0]) #tanggal
        slide1 = prs.slides.add_slide(prs.slide_layouts[1]) #sumatra
        slide2 = prs.slides.add_slide(prs.slide_layouts[1]) #dewata
        slide3 = prs.slides.add_slide(prs.slide_layouts[1]) #karimun
        slide4 = prs.slides.add_slide(prs.slide_layouts[1]) #derawan
        slide5 = prs.slides.add_slide(prs.slide_layouts[1]) #oceanflow1
        slide19 = prs.slides.add_slide(prs.slide_layouts[1]) #celebes
        slide20 = prs.slides.add_slide(prs.slide_layouts[1]) #greencalypso
        slide6 = prs.slides.add_slide(prs.slide_layouts[1]) #sumba
        slide22 = prs.slides.add_slide(prs.slide_layouts[1]) #putrialysha
        slide24 = prs.slides.add_slide(prs.slide_layouts[1]) #bunaken
        slide17 = prs.slides.add_slide(prs.slide_layouts[1]) #java
        slide18 = prs.slides.add_slide(prs.slide_layouts[1]) #natuna
        slide26 = prs.slides.add_slide(prs.slide_layouts[1]) #svii


        slide7 = prs.slides.add_slide(prs.slide_layouts[3]) #sumatra
        slide8 = prs.slides.add_slide(prs.slide_layouts[3]) #dewata
        slide9 = prs.slides.add_slide(prs.slide_layouts[3]) #karimun
        slide10 = prs.slides.add_slide(prs.slide_layouts[3]) #derawan
        slide11 = prs.slides.add_slide(prs.slide_layouts[3]) #oceanflow1
        slide12 = prs.slides.add_slide(prs.slide_layouts[3]) #sumba
        slide13 = prs.slides.add_slide(prs.slide_layouts[3]) #java
        slide14 = prs.slides.add_slide(prs.slide_layouts[3]) #natuna
        slide27 = prs.slides.add_slide(prs.slide_layouts[3]) #svii
        slide15 = prs.slides.add_slide(prs.slide_layouts[3]) #celebes
        slide21 = prs.slides.add_slide(prs.slide_layouts[3]) #greencalypso
        slide23 = prs.slides.add_slide(prs.slide_layouts[3]) #putrialysha
        slide25 = prs.slides.add_slide(prs.slide_layouts[3]) #bunaken
        slide16 = prs.slides.add_slide(prs.slide_layouts[3]) #borneo

        # title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Operation Update"
        subtitle.text = str(dat_sumatra['Month'].iloc[-2])


        # slide1 Berau Coal
        slide1.placeholders[16].text = 'BULK SUMATRA \nBERAU COAL OPERATIONAL DASHBOARD'
        slide1.placeholders[17].text = 'Berau Coal' 
        slide1.shapes.add_picture('img/NLR_type Bulk Sumatra.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide1.shapes.add_picture('img/GLR_type Bulk Sumatra.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide1.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide1.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))


        # slide2 Berau Coal
        slide2.placeholders[16].text = 'BULK DEWATA \nBERAU COAL OPERATIONAL DASHBOARD'
        slide2.placeholders[17].text = 'Berau Coal' 
        slide2.shapes.add_picture('img/NLR_type Bulk Dewata.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide2.shapes.add_picture('img/GLR_type Bulk Dewata.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide2.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide2.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide3 Berau Coal
        slide3.placeholders[16].text = 'BULK KARIMUN \nBERAU COAL OPERATIONAL DASHBOARD'
        slide3.placeholders[17].text = 'Berau Coal'
        slide3.shapes.add_picture('img/NLR_type Bulk Karimun.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide3.shapes.add_picture('img/GLR_type Bulk Karimun.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide3.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide3.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide4 Berau Coal
        slide4.placeholders[16].text = 'BULK DERAWAN \nBERAU COAL OPERATIONAL DASHBOARD'
        slide4.placeholders[17].text = 'Berau Coal'
        slide4.shapes.add_picture('img/NLR_type Bulk Derawan.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide4.shapes.add_picture('img/GLR_type Bulk Derawan.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide4.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide4.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide5 Berau Coal
        slide5.placeholders[16].text = 'BULK OCEAN FLOW 1 \nBERAU COAL OPERATIONAL DASHBOARD'
        slide5.placeholders[17].text = 'Berau Coal'
        slide5.shapes.add_picture('img/NLR_type Ocean Flow 1.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide5.shapes.add_picture('img/GLR_type Ocean Flow 1.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide5.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide5.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide19 Berau Coal
        slide19.placeholders[16].text = 'BULK CELEBES \nBERAU COAL OPERATIONAL DASHBOARD'
        slide19.placeholders[17].text = 'Berau Coal'
        slide19.shapes.add_picture('img/NLR_type Bulk Celebes.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide19.shapes.add_picture('img/GLR_type Bulk Celebes.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide19.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide19.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide20 Berau Coal
        slide20.placeholders[16].text = 'GREEN CALYPSO \nBERAU COAL OPERATIONAL DASHBOARD'
        slide20.placeholders[17].text = 'Berau Coal'
        slide20.shapes.add_picture('img/NLR_type Green Calypso.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide20.shapes.add_picture('img/GLR_type Green Calypso.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide20.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide20.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide6 BGE
        slide6.placeholders[16].text = 'BULK SUMBA \nBGE OPERATIONAL DASHBOARD'
        slide6.placeholders[17].text = 'BGE'
        slide6.shapes.add_picture('img/NLR_type Bulk Sumba.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide6.shapes.add_picture('img/GLR_type Bulk Sumba.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide6.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide6.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide22 BGE
        slide22.placeholders[16].text = 'PUTRI ALYSHA \nBGE OPERATIONAL DASHBOARD'
        slide22.placeholders[17].text = 'BGE'
        slide22.shapes.add_picture('img/NLR_type Putri Alysha.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide22.shapes.add_picture('img/GLR_type Putri Alysha.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide22.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide22.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide24 BGE
        slide24.placeholders[16].text = 'BULK BUNAKEN \nBGE OPERATIONAL DASHBOARD'
        slide24.placeholders[17].text = 'BGE'
        slide24.shapes.add_picture('img/NLR_type Bulk Bunaken.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide24.shapes.add_picture('img/GLR_type Bulk Bunaken.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide24.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide24.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide17 PSS
        slide17.placeholders[16].text = 'BULK JAVA \nPSS OPERATIONAL DASHBOARD'
        slide17.placeholders[17].text = 'PSS'
        slide17.shapes.add_picture('img/NLR_type Bulk Java.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide17.shapes.add_picture('img/GLR_type Bulk Java.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide17.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide17.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide18 BIB
        slide18.placeholders[16].text = 'BULK NATUNA \nBIB OPERATIONAL DASHBOARD'
        slide18.placeholders[17].text = 'BIB'
        slide18.shapes.add_picture('img/NLR_type Bulk Natuna.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide18.shapes.add_picture('img/GLR_type Bulk Natuna.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide18.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide18.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide26 BIB
        slide26.placeholders[16].text = 'STRAITS VENTURE II \nBIB OPERATIONAL DASHBOARD'
        slide26.placeholders[17].text = 'BIB'
        slide26.shapes.add_picture('img/NLR_type STRAITS VENTURE II.png', Cm(1.5), Cm(3.5), Cm(15.2), Cm(10))
        slide26.shapes.add_picture('img/GLR_type STRAITS VENTURE II.png', Cm(17.5), Cm(3.5), Cm(15.2), Cm(10))

        slide26.shapes.add_picture('img/yellow bullet.png', Cm(1.75), Cm(14))
        slide26.shapes.add_picture('img/yellow bullet.png', Cm(17.75), Cm(14))

        # slide7 Bulk Sumatra
        slide7.placeholders[10].text = 'BULK SUMATRA \nOPS PERFORMANCE DASHBOARD'
        slide7.placeholders[17].text = 'Berau Coal'
        slide7.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide7.shapes.add_picture('img/Volume Bulk Sumatra.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide7.shapes.add_picture('img/Fuel Ratio Bulk Sumatra.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide7.shapes.add_picture('img/GLR Bulk Sumatra.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide7.shapes.add_picture('img/NLR Bulk Sumatra.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide8 Bulk Dewata
        slide8.placeholders[10].text = 'BULK DEWATA \nOPS PERFORMANCE DASHBOARD'
        slide8.placeholders[17].text = 'Berau Coal'
        slide8.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide8.shapes.add_picture('img/Volume Bulk Dewata.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide8.shapes.add_picture('img/Fuel Ratio Bulk Dewata.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide8.shapes.add_picture('img/GLR Bulk Dewata.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide8.shapes.add_picture('img/NLR Bulk Dewata.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide9 Bulk Karimun
        slide9.placeholders[10].text = 'BULK KARIMUN \nOPS PERFORMANCE DASHBOARD'
        slide9.placeholders[17].text = 'Berau Coal'
        slide9.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide9.shapes.add_picture('img/Volume Bulk Karimun.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide9.shapes.add_picture('img/Fuel Ratio Bulk Karimun.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide9.shapes.add_picture('img/GLR Bulk Karimun.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide9.shapes.add_picture('img/NLR Bulk Karimun.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide10 Bulk Derawan
        slide10.placeholders[10].text = 'BULK DERAWAN \nOPS PERFORMANCE DASHBOARD'
        slide10.placeholders[17].text = 'Berau Coal'
        slide10.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide10.shapes.add_picture('img/Volume Bulk Derawan.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide10.shapes.add_picture('img/Fuel Ratio Bulk Derawan.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide10.shapes.add_picture('img/GLR Bulk Derawan.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide10.shapes.add_picture('img/NLR Bulk Derawan.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide11 Ocean Flow 1
        slide11.placeholders[10].text = 'OCEAN FLOW 1 \nOPS PERFORMANCE DASHBOARD'
        slide11.placeholders[17].text = 'Berau Coal'
        slide11.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide11.shapes.add_picture('img/Volume Ocean Flow 1.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide11.shapes.add_picture('img/Fuel Ratio Ocean Flow 1.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide11.shapes.add_picture('img/GLR Ocean Flow 1.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide11.shapes.add_picture('img/NLR Ocean Flow 1.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide12 Bulk Sumba
        slide12.placeholders[10].text = 'BULK SUMBA \nOPS PERFORMANCE DASHBOARD'
        slide12.placeholders[17].text = 'BGE'
        slide12.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide12.shapes.add_picture('img/Volume Bulk Sumba.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide12.shapes.add_picture('img/Fuel Ratio Bulk Sumba.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide12.shapes.add_picture('img/GLR Bulk Sumba.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide12.shapes.add_picture('img/NLR Bulk Sumba.png', Cm(14), Cm(10.3), Cm(12), Cm(7))
                            
        # slide13 Bulk Java
        slide13.placeholders[10].text = 'BULK JAVA \nOPS PERFORMANCE DASHBOARD'
        slide13.placeholders[17].text = 'PSS'
        slide13.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide13.shapes.add_picture('img/Volume Bulk Java.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide13.shapes.add_picture('img/Fuel Ratio Bulk Java.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide13.shapes.add_picture('img/GLR Bulk Java.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide13.shapes.add_picture('img/NLR Bulk Java.png', Cm(14), Cm(10.3), Cm(12), Cm(7))
                            
        # slide14 Bulk Natuna
        slide14.placeholders[10].text = 'BULK NATUNA \nOPS PERFORMANCE DASHBOARD'
        slide14.placeholders[17].text = 'BIB'
        slide14.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide14.shapes.add_picture('img/Volume Bulk Natuna.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide14.shapes.add_picture('img/Fuel Ratio Bulk Natuna.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide14.shapes.add_picture('img/GLR Bulk Natuna.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide14.shapes.add_picture('img/NLR Bulk Natuna.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide27 STRAITS VENTURE II
        slide27.placeholders[10].text = 'STRAITS VENTURE II \nOPS PERFORMANCE DASHBOARD'
        slide27.placeholders[17].text = 'BIB'
        slide27.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide27.shapes.add_picture('img/Volume STRAITS VENTURE II.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide27.shapes.add_picture('img/Fuel Ratio STRAITS VENTURE II.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide27.shapes.add_picture('img/GLR STRAITS VENTURE II.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide27.shapes.add_picture('img/NLR STRAITS VENTURE II.png', Cm(14), Cm(10.3), Cm(12), Cm(7))
                            
        # slide15 Bulk Celebes
        slide15.placeholders[10].text = 'BULK CELEBES \nOPS PERFORMANCE DASHBOARD'
        slide15.placeholders[17].text = 'Berau Coal'
        slide15.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide15.shapes.add_picture('img/Volume Bulk Celebes.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide15.shapes.add_picture('img/Fuel Ratio Bulk Celebes.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide15.shapes.add_picture('img/GLR Bulk Celebes.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide15.shapes.add_picture('img/NLR Bulk Celebes.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide21 Green Calypso
        slide21.placeholders[10].text = 'GREEN CALYPSO \nOPS PERFORMANCE DASHBOARD'
        slide21.placeholders[17].text = 'Berau Coal'
        slide21.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide21.shapes.add_picture('img/Volume Green Calypso.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide21.shapes.add_picture('img/Fuel Ratio Green Calypso.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide21.shapes.add_picture('img/GLR Green Calypso.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide21.shapes.add_picture('img/NLR Green Calypso.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide23 Putri Alysha
        slide23.placeholders[10].text = 'PUTRI ALYSHA \nOPS PERFORMANCE DASHBOARD'
        slide23.placeholders[17].text = 'BGE'
        slide23.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide23.shapes.add_picture('img/Volume Putri Alysha.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide23.shapes.add_picture('img/Fuel Ratio Putri Alysha.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide23.shapes.add_picture('img/GLR Putri Alysha.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide23.shapes.add_picture('img/NLR Putri Alysha.png', Cm(14), Cm(10.3), Cm(12), Cm(7))

        # slide25 Bulk Bunaken
        slide25.placeholders[10].text = 'BULK BUNAKEN \nOPS PERFORMANCE DASHBOARD'
        slide25.placeholders[17].text = 'BGE'
        slide25.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide25.shapes.add_picture('img/Volume Bulk Bunaken.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide25.shapes.add_picture('img/Fuel Ratio Bulk Bunaken.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide25.shapes.add_picture('img/GLR Bulk Bunaken.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide25.shapes.add_picture('img/NLR Bulk Bunaken.png', Cm(14), Cm(10.3), Cm(12), Cm(7))
                            
        # slide16 Bulk Borneo
        slide16.placeholders[10].text = 'BULK BORNEO \nOPS PERFORMANCE DASHBOARD'
        slide16.placeholders[17].text = 'LDPL-KAMSAR'
        slide16.placeholders[11].text = 'As of ' + str(date.today().strftime("%d %b %Y"))
        slide16.shapes.add_picture('img/Volume Bulk Borneo.png', Cm(1.25), Cm(3), Cm(12), Cm(7))
        slide16.shapes.add_picture('img/Fuel Ratio Bulk Borneo.png', Cm(14), Cm(3),Cm(12), Cm(7))
        slide16.shapes.add_picture('img/GLR Bulk Borneo.png', Cm(1.25), Cm(10.3), Cm(12), Cm(7))
        slide16.shapes.add_picture('img/NLR Bulk Borneo.png', Cm(14), Cm(10.3), Cm(12), Cm(7))                        


        # Saving the PowerPoint presentation
        prs.save(bytes_io)    
    
    return dcc.send_bytes(to_pptx, 'CTS Performance Dashboard.pptx'), download_date


######################
# Plot Function
######################
##-----Function
#-- 1. Volume Function
def plot_volume(df, title):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['Volume Plan'].tail(4),
        name='Plan',
        text=df['Volume Plan'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#5b9bd5'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['Volume Actual'].tail(4),
        name='Actual',
        text=df['Volume Actual'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='lightsalmon'))
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.25,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':40}},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['Volume Plan'].iloc[-1])*1.8],
        # yaxis_range=[0,850],
        bargroupgap=0.165,
        bargap=0.25,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.25,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])

    labels = list(df['Month'])
    labels[-1] = "YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 2. NLR Function
def plot_nlr(df, title):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Plan'].tail(4),
        name='Plan',
        text=df['NLR Plan'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#5b9bd5'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Actual'].tail(4),
        name='Actual',
        text=df['NLR Actual'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='lightsalmon'))
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.25,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':40}},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['NLR Plan'].iloc[-1])*1.65],
        bargroupgap=0.165,
        bargap=0.25,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.25,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])

    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 3. GLR Function
def plot_glr(df, title):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Plan'].tail(4),
        name='Plan',
        text=df['GLR Plan'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#5b9bd5'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Actual'].tail(4),
        name='Actual',
        text=df['GLR Actual'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='lightsalmon'))
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.25,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':40}},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['GLR Plan'].iloc[-1])*1.65],
        bargroupgap=0.165,
        bargap=0.25,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.25,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])

    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 4. Fuel Ratio Function
def plot_fr(df, title, baseline):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['Fuel Ratio Gross'].tail(4),
        name='Gross',
        text=df['Fuel Ratio Gross'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#5b9bd5'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['Fuel Ratio Net'].tail(4),
        name='Net',
        text=df['Fuel Ratio Net'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='lightsalmon'))
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.25,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':40}},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['Fuel Ratio Gross'].max())*1.65],
        bargroupgap=0.165,
        bargap=0.25,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.25,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])

    fig.add_hline(y=baseline, line_width=2, line_dash="dash", line_color="red")
    fig.add_annotation(
        x=3.5,
        y=baseline+0.03,
        xref="x",
        yref="y",
        text=baseline,
        showarrow=False,
        font=dict(
            family="Verdana",
            size=48,
            color="#ffffff"
            ),
        align="center",
        bordercolor="#c7c7c7",
        borderwidth=2,
        borderpad=4,
        bgcolor="red",
        opacity=0.8
        )
    
    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 5. NLR Type Function
def plot_nlr_type(df, title, baseline):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Single'].tail(4),
        name='Single',
        text=df['NLR Single'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#ffc000'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Blending'].tail(4),
        name='Blending',
        text=df['NLR Blending'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#00b1f1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Gear'].tail(4),
        name='Gear',
        text=df['NLR Gear'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#0071c1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Barge'].tail(4),
        name='Barge',
        text=df['NLR Barge'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#92d14f'))
    
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.2,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':30},
            'orientation':'h'},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['NLR Plan'].iloc[-1])*1.65],
        bargroupgap=0.05,
        bargap=0.15,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.2,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])
    
    fig.add_hline(y=baseline, line_width=1, line_dash="dash", line_color="red")
    fig.add_annotation(
        x=3.75,
        y=baseline*1.07,
        xref="x",
        yref="y",
        text=str(baseline) + ' MT/Day',
        showarrow=False,
        font=dict(
            family="Verdana",
            size=27,
            color="#ffffff"
            ),
        align="center",
        bordercolor="#c7c7c7",
        borderwidth=2,
        borderpad=4,
        bgcolor="red",
        opacity=0.8
        )

    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 5B. NLR Type Function Double Baseline
def plot_nlr_type2(df, title, baseline1, baseline2):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Single'].tail(4),
        name='Single',
        text=df['NLR Single'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#ffc000'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Blending'].tail(4),
        name='Blending',
        text=df['NLR Blending'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#00b1f1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Gear'].tail(4),
        name='Gear',
        text=df['NLR Gear'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#0071c1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['NLR Barge'].tail(4),
        name='Barge',
        text=df['NLR Barge'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#92d14f'))
    
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.2,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':30},
            'orientation':'h'},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['NLR Plan'].iloc[-1])*1.65],
        bargroupgap=0.05,
        bargap=0.15,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.2,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])
    
    fig.add_hline(y=baseline1, line_width=1, line_dash="dash", line_color="red")
    fig.add_annotation(
        x=3.75,
        y=baseline1*1.07,
        xref="x",
        yref="y",
        text=str(baseline1) + ' MT/Day',
        showarrow=False,
        font=dict(
            family="Verdana",
            size=27,
            color="#ffffff"
            ),
        align="center",
        bordercolor="#c7c7c7",
        borderwidth=2,
        borderpad=4,
        bgcolor="red",
        opacity=0.8
        )
    
    fig.add_hline(y=baseline2, line_width=1, line_dash="dash", line_color="red")
    fig.add_annotation(
        x=3.75,
        y=baseline2*1.07,
        xref="x",
        yref="y",
        text=str(baseline2) + ' MT/Day',
        showarrow=False,
        font=dict(
            family="Verdana",
            size=27,
            color="#ff7f7f"
            ),
        align="center",
        bordercolor="#c7c7c7",
        borderwidth=2,
        borderpad=4,
        bgcolor="red",
        opacity=0.8
        )

    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)

#-- 6. GLR Type Function
def plot_glr_type(df, title):
    fig = go.Figure()
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Single'].tail(4),
        name='Single',
        text=df['GLR Single'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#ffc000'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Blending'].tail(4),
        name='Blending',
        text=df['GLR Blending'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#00b1f1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Gear'].tail(4),
        name='Gear',
        text=df['GLR Gear'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#0071c1'))
    fig.add_trace(go.Bar(
        x=df['Month'].tail(4),
        y=df['GLR Barge'].tail(4),
        name='Barge',
        text=df['GLR Barge'].tail(4),
        textfont_size=48,
        textposition = 'outside',
        hovertemplate='%{y:y}',
        textangle=0,
        marker_color='#92d14f'))
    
    fig.update_layout({
        'height':1000,'width':1700,
        'margin' : {'t':200, 'b':3, 'l':3, 'r':3},
        "autosize": True,
        'plot_bgcolor': 'rgba(0, 0, 0, 0)',
        'paper_bgcolor': 'rgba(0, 0, 0, 0)',},
        title={
            'text': title,
            'y':0.89,
            'x':0.5,
            'xanchor': 'center',
            'yanchor': 'top',
            'font': {'size':75}},
        legend={
            'yanchor':"bottom",
            'y':-0.2,
            'xanchor':"center",
            'x':0.5,
            'itemsizing': 'constant',
            'font':{'size':30},
            'orientation':'h'},
        xaxis = dict(tickfont = dict(size=40)),
        yaxis = dict(tickfont = dict(size=40)),
        yaxis_range=[0,(df['GLR Plan'].iloc[-1])*1.65],
        bargroupgap=0.05,
        bargap=0.15,
        hovermode="x",
        shapes=[go.layout.Shape(type='rect', 
                                xref='paper',
                                yref='paper',
                                x0=0,
                                y0=-0.2,
                                x1=1,
                                y1=1.25,
                                line={'width': 2, 'color': 'black'})])

    labels = list(df['Month'])
    labels[-1] = "AVG YTD'25"

    fig.update_xaxes(linecolor='black', tickvals=np.arange(4), ticktext=labels[-4:])
    fig.update_yaxes(showgrid=False, visible=False)
    
    return(fig)